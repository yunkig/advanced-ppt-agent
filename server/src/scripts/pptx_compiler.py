#!/usr/bin/env python3
"""
Python PPTX 真实二进制编译器
读取 Sub-Agent 协同生成的结构化 DSL JSON，直接调用 python-pptx 编译生成可供本地 PowerPoint/WPS 任意二次编辑的原生 .pptx 文件。
"""

import sys
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    if len(hex_str) == 3:
        hex_str = ''.join([c*2 for c in hex_str])
    return RGBColor(*(int(hex_str[i:i+2], 16) for i in (0, 2, 4)))

def compile_pptx(json_path, output_path):
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    plan = data.get('plan', {})
    template = data.get('template', {})
    slides_data = data.get('slides', [])
    palette = template.get('palette', {})

    primary_color = hex_to_rgb(palette.get('primary', '#0F172A'))
    secondary_color = hex_to_rgb(palette.get('secondary', '#1E293B'))
    accent_color = hex_to_rgb(palette.get('accent', '#3B82F6'))
    bg_color = hex_to_rgb(palette.get('background', '#F8FAFC'))

    prs = Presentation()
    # 设置 16:9 黄金宽屏
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # 1. 封面页
    slide_cover = prs.slides.add_slide(blank_layout)
    
    # 封面背景块
    bg_shape = slide_cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = primary_color
    bg_shape.line.color.rgb = primary_color

    # 装饰线条
    deco_line = slide_cover.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(2.2), Inches(0.15), Inches(2.8))
    deco_line.fill.solid()
    deco_line.fill.fore_color.rgb = accent_color
    deco_line.line.color.rgb = accent_color

    # 封面大标题
    txBox = slide_cover.shapes.add_textbox(Inches(1.6), Inches(2.2), Inches(10.5), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = plan.get('topic', '企业级智能演示文稿')
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 封面副标题
    p2 = tf.add_paragraph()
    p2.text = f"基于 Multi-Agent 架构的端到端高质量方案 · 面向受众：{plan.get('targetAudience', '专业团队')}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(200, 210, 230)
    p2.space_before = Pt(20)

    # 2. 依次编译各个 Sub-Agent 交付的幻灯片页面
    for idx, s in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        
        # 页面底色背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.color.rgb = bg_color

        # 顶部导航色带
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = primary_color
        top_bar.line.color.rgb = primary_color

        # 页面标题
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(10.0), Inches(0.8))
        t_frame = t_box.text_frame
        t_p = t_frame.paragraphs[0]
        sub_id = s.get('subId', '')
        t_p.text = f"{sub_id}  {s.get('title', '核心论述')}" if sub_id else s.get('title', '核心论述')
        t_p.font.size = Pt(22)
        t_p.font.bold = True
        t_p.font.color.rgb = RGBColor(255, 255, 255)

        # 页码指示器
        page_box = slide.shapes.add_textbox(Inches(11.2), Inches(0.3), Inches(1.5), Inches(0.5))
        page_p = page_box.text_frame.paragraphs[0]
        page_p.text = f"{idx + 2} / {len(slides_data) + 1}"
        page_p.font.size = Pt(13)
        page_p.font.color.rgb = RGBColor(200, 210, 230)
        page_p.alignment = PP_ALIGN.RIGHT

        # 判断是否为图表页
        chart_data_info = s.get('chartData')
        if s.get('pageType') == 'chart' and chart_data_info:
            # 左侧文字卡片 (Inches 5.5 宽)
            left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.2), Inches(5.2))
            left_card.fill.solid()
            left_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            left_card.line.color.rgb = RGBColor(220, 225, 235)

            lt_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(4.8), Inches(4.8))
            lt_frame = lt_box.text_frame
            lt_frame.word_wrap = True
            lt_head = lt_frame.paragraphs[0]
            lt_head.text = "指标解读与核心洞察"
            lt_head.font.size = Pt(16)
            lt_head.font.bold = True
            lt_head.font.color.rgb = primary_color

            for bp in s.get('bulletPoints', []):
                p_bp = lt_frame.add_paragraph()
                p_bp.text = f"• {bp}"
                p_bp.font.size = Pt(13)
                p_bp.font.color.rgb = RGBColor(60, 70, 85)
                p_bp.space_before = Pt(10)

            # 右侧原生 Excel 图表
            chart_data = CategoryChartData()
            chart_data.categories = chart_data_info.get('categories', ['Q1', 'Q2', 'Q3', 'Q4'])
            series_list = chart_data_info.get('series', [{'name': '增长效能(%)', 'values': [35, 58, 72, 84]}])
            for ser in series_list:
                chart_data.add_series(ser.get('name', '指标'), ser.get('values', [20, 40, 60, 80]))

            x, y, cx, cy = Inches(6.4), Inches(1.6), Inches(6.1), Inches(5.2)
            chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
            chart.has_legend = True
            chart.legend.include_in_layout = False
        else:
            # 常规两栏/三栏内容卡片
            bullet_points = s.get('bulletPoints', ['核心业务要点', '落地执行策略', '预期达成成效'])
            card_count = min(len(bullet_points), 3)
            if card_count == 0:
                card_count = 1
                bullet_points = ['详实阐述与数据依据']

            card_width = (Inches(11.7) - (Inches(0.4) * (card_count - 1))) / card_count
            for c_idx in range(card_count):
                c_left = Inches(0.8) + c_idx * (card_width + Inches(0.4))
                
                card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, Inches(1.6), card_width, Inches(5.2))
                card_shape.fill.solid()
                card_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                card_shape.line.color.rgb = RGBColor(220, 225, 235)

                # 顶部强调条
                strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, c_left, Inches(1.6), card_width, Inches(0.12))
                strip.fill.solid()
                strip.fill.fore_color.rgb = accent_color
                strip.line.color.rgb = accent_color

                c_box = slide.shapes.add_textbox(c_left + Inches(0.25), Inches(1.9), card_width - Inches(0.5), Inches(4.6))
                c_frame = c_box.text_frame
                c_frame.word_wrap = True
                
                c_head = c_frame.paragraphs[0]
                c_head.text = f"核心维度 0{c_idx + 1}"
                c_head.font.size = Pt(15)
                c_head.font.bold = True
                c_head.font.color.rgb = primary_color

                c_desc = c_frame.add_paragraph()
                c_desc.text = bullet_points[c_idx] if c_idx < len(bullet_points) else "支撑论述与量化依据"
                c_desc.font.size = Pt(13)
                c_desc.font.color.rgb = RGBColor(60, 70, 85)
                c_desc.space_before = Pt(14)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    print(f"SUCCESS: Compiled native presentation to {output_path}")

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python3 pptx_compiler.py <input_json_path> <output_pptx_path>")
        sys.exit(1)
    compile_pptx(sys.argv[1], sys.argv[2])
